from PyPDF2 import PdfReader
import docx
from pptx import Presentation


def lire_pdf(fichier):
    lecteur_pdf = PdfReader(fichier)
    texte = ""
    for page in lecteur_pdf.pages:
        texte += page.extract_text()
    return texte

def lire_docx(fichier):
    document = docx.Document(fichier)
    texte = ""
    for paragraphe in document.paragraphs:
        texte += paragraphe.text + "\n"
    return texte

def lire_ppt(fichier):
    presentation = Presentation(fichier)
    texte = ""
    for slide in presentation.slides:
        for element in slide.shapes:
            if hasattr(element, "text"):
                texte += element.text + "\n"
    return texte

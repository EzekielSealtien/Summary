import streamlit as st
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

def lire_pdf(fichier):
    lecteur_pdf = PdfReader(fichier)
    texte = ""
    for page in lecteur_pdf.pages:
        texte += page.extract_text()
    return texte

def lire_docx(fichier):
    document = docx.Document(fichier)
    texte = ""
    for paragraphe in document.paragraphs:
        texte += paragraphe.text + "\n"
    return texte

def lire_ppt(fichier):
    presentation = Presentation(fichier)
    texte = ""
    for slide in presentation.slides:
        for element in slide.shapes:
            if hasattr(element, "text"):
                texte += element.text + "\n"
    return texte

st.markdown(
    """
    <style>
    .stApp {
        background-color: #d0f0c0;
    }
    </style>
    """, 
    unsafe_allow_html=True
)

st.markdown("<h1 style='text-align: center; color: #4CAF50;'>📄 Uploader et Afficher le Contenu de votre Fichier</h1>", unsafe_allow_html=True)

st.markdown("---")
st.markdown("<h4 style='text-align: center;'>Sélectionnez un fichier pour afficher son contenu</h4>", unsafe_allow_html=True)

fichier_telecharge = st.file_uploader("", type=["pdf", "docx", "pptx"])


file_content=""

if fichier_telecharge is not None:
    extension_fichier = fichier_telecharge.name.split('.')[-1].lower()

    if extension_fichier in ["pdf", "docx", "pptx"]:
        st.markdown("---")
        st.markdown(f"<h5 style='color: #FF5733;'>Fichier détecté : {fichier_telecharge.name}</h5>", unsafe_allow_html=True)
        
        if extension_fichier == "pdf":
            file_content = lire_pdf(fichier_telecharge)
        elif extension_fichier == "docx":
            file_content = lire_docx(fichier_telecharge)
        elif extension_fichier == "pptx":
            file_content = lire_ppt(fichier_telecharge)
            
    else:
        st.markdown("---")
        st.error("⚠️ Seuls les fichiers PDF, DOCX, et PPTX sont acceptés.")

#Invocation du modèle---> LLM Amazon Tittan 




# Lorsqu'on clique sur le bouton Afficher le resumé
if st.button("Afficher le resumé :"):
    st.write(file_content)

st.markdown("---")
st.markdown("<p style='text-align: center; color: grey;'>App conçu l'equipe RAD</p>", unsafe_allow_html=True)
